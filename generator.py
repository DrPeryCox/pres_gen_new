# generator.py
import base64
import io
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt, Inches
from pptx.slide import Slide

# Константы
LAYOUT_TITLE_ONLY = 5
LAYOUT_TITLE_AND_CONTENT = 1
LAYOUT_TWO_CONTENT = 3
PLACEHOLDER_IDX_TITLE = 0
PLACEHOLDER_IDX_CONTENT_CENTER = 1
PLACEHOLDER_IDX_CONTENT_LEFT = 1
PLACEHOLDER_IDX_CONTENT_RIGHT = 2
CONTENT_TYPE_TEXT = "content"
CONTENT_TYPE_BULLETS = "bullet_points"
CONTENT_TYPE_IMAGE = "image"
CONTENT_TYPE_BULLETS_HEADER = "bullet_points_header"


class PresentationGenerator:
    """
    Генерирует .pptx файл на основе словаря, полученного из JSON.
    """

    def __init__(self, data: Dict[str, Any]):
        if "slides" not in data:
            raise KeyError("Ключ 'slides' не найден в предоставленных данных.")
        self.data = data
        self.prs = Presentation()
        self.prs.slide_width = Inches(10.8)
        self.prs.slide_height = Inches(10.8)

    def generate(self) -> io.BytesIO:
        print("Начинаю генерацию презентации...")
        for i, slide_data in enumerate(self.data.get("slides", [])):
            print(f"  - Обработка слайда {i + 1}...")
            try:
                self._add_slide(slide_data)
            except (KeyError, ValueError) as e:
                print(f"    [ОШИБКА] Пропуск слайда {i + 1} из-за ошибки: {e}")
                raise ValueError(f"Ошибка на слайде {i + 1}: {e}") from e

        file_stream = io.BytesIO()
        self.prs.save(file_stream)
        file_stream.seek(0)
        print("Презентация успешно сгенерирована в памяти.")
        return file_stream

    @staticmethod
    def _is_part_meaningful(part_data: Optional[Dict[str, Any]]) -> bool:
        if not part_data:
            return False
        return bool(part_data.get(CONTENT_TYPE_TEXT) or
                    part_data.get(CONTENT_TYPE_BULLETS) or
                    part_data.get(CONTENT_TYPE_IMAGE))

    def _add_slide(self, slide_data: Dict[str, Any]) -> None:
        layout = self._get_layout(slide_data)
        slide = self.prs.slides.add_slide(layout)

        self._set_title(slide, slide_data)
        self._set_background(slide, slide_data.get("background"))

        if self._is_part_meaningful(slide_data.get("center_part")):
            placeholder = slide.placeholders[PLACEHOLDER_IDX_CONTENT_CENTER]
            self._process_content_part(slide, placeholder, slide_data["center_part"])

        elif self._is_part_meaningful(slide_data.get("left_part")) or \
                self._is_part_meaningful(slide_data.get("right_part")):

            if self._is_part_meaningful(slide_data.get("left_part")):
                placeholder_left = slide.placeholders[PLACEHOLDER_IDX_CONTENT_LEFT]
                self._process_content_part(slide, placeholder_left, slide_data["left_part"])

            if self._is_part_meaningful(slide_data.get("right_part")):
                placeholder_right = slide.placeholders[PLACEHOLDER_IDX_CONTENT_RIGHT]
                self._process_content_part(slide, placeholder_right, slide_data["right_part"])

    def _get_layout(self, slide_data: Dict[str, Any]):
        has_center = self._is_part_meaningful(slide_data.get("center_part"))
        has_left = self._is_part_meaningful(slide_data.get("left_part"))
        has_right = self._is_part_meaningful(slide_data.get("right_part"))

        if has_center:
            return self.prs.slide_layouts[LAYOUT_TITLE_AND_CONTENT]
        if has_left or has_right:
            return self.prs.slide_layouts[LAYOUT_TWO_CONTENT]

        return self.prs.slide_layouts[LAYOUT_TITLE_ONLY]

    def _set_title(self, slide: Slide, slide_data: Dict[str, Any]) -> None:
        if not slide.shapes.title: return
        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", " ")
        font = title_shape.text_frame.paragraphs[0].font
        if "font_size" in slide_data: font.size = Pt(slide_data["font_size"])
        if "font_color" in slide_data: font.color.rgb = RGBColor(*slide_data["font_color"])

    def _set_background(self, slide: Slide, b64_image: Optional[str]) -> None:
        if not b64_image: return
        try:
            image_stream = self._decode_base64_to_stream(b64_image)
            slide.background.fill.solid()
            slide.background.fill.picture(image_stream)
        except Exception as e:
            print(f"    [ПРЕДУПРЕЖДЕНИЕ] Не удалось установить фон: {e}")

    def _process_content_part(self, slide: Slide, placeholder, part_data: Dict[str, Any]) -> None:
        content_keys = [key for key in [CONTENT_TYPE_TEXT, CONTENT_TYPE_BULLETS, CONTENT_TYPE_IMAGE] if
                        key in part_data and part_data[key]]
        if len(content_keys) > 1: raise ValueError(
            f"Часть контента может содержать только один из ключей: {', '.join(content_keys)}")
        if not content_keys: return

        content_type = content_keys[0]
        font_size = part_data.get("font_size")
        font_color = part_data.get("font_color")

        if content_type == CONTENT_TYPE_TEXT:
            self._add_text(slide, placeholder, part_data[CONTENT_TYPE_TEXT], font_size, font_color)
        elif content_type == CONTENT_TYPE_BULLETS:
            header = part_data.get(CONTENT_TYPE_BULLETS_HEADER)
            self._add_bullet_points(slide, placeholder, part_data[CONTENT_TYPE_BULLETS], header, font_size, font_color)
        elif content_type == CONTENT_TYPE_IMAGE:
            self._add_image(slide, placeholder, part_data[CONTENT_TYPE_IMAGE])

    def _add_text(self, slide: Slide, placeholder, text: str, font_size: Optional[int],
                  font_color: Optional[List[int]]) -> None:
        left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame

        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = text
        self._apply_font_style(p.font, font_size, font_color)

        sp = placeholder.element
        sp.getparent().remove(sp)

    def _add_bullet_points(self, slide: Slide, placeholder, points: List[str], header: Optional[str],
                           font_size: Optional[int],
                           font_color: Optional[List[int]]) -> None:
        if not header and not points:
            sp = placeholder.element
            sp.getparent().remove(sp)
            return

        if not header:
            text_frame = placeholder.text_frame
            text_frame.clear()
            for i, point_text in enumerate(points):
                p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
                p.text = point_text
                p.level = 0
                self._apply_font_style(p.font, font_size, font_color)
            return

        left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height

        header_height = Pt(40)
        txBox_header = slide.shapes.add_textbox(left, top, width, header_height)
        p_header = txBox_header.text_frame.paragraphs[0]
        p_header.text = header
        self._apply_font_style(p_header.font, font_size, font_color)
        p_header.font.bold = True

        list_top = top + header_height
        list_height = height - header_height

        if list_height > Pt(20) and points:
            txBox_list = slide.shapes.add_textbox(left, list_top, width, list_height)
            tf_list = txBox_list.text_frame

            tf_list.auto_size = MSO_AUTO_SIZE.NONE
            tf_list.word_wrap = True

            tf_list.clear()

            for i, point_text in enumerate(points):
                p_list = tf_list.paragraphs[0] if i == 0 else tf_list.add_paragraph()
                p_list.text = f"• {point_text}"
                p_list.level = 0
                self._apply_font_style(p_list.font, font_size, font_color)

        sp = placeholder.element
        sp.getparent().remove(sp)

    def _add_image(self, slide: Slide, placeholder, b64_image: str) -> None:
        if not b64_image: return
        try:
            image_stream = self._decode_base64_to_stream(b64_image)
            slide.shapes.add_picture(image_stream, placeholder.left, placeholder.top, width=placeholder.width,
                                     height=placeholder.height)
            sp = placeholder.element
            sp.getparent().remove(sp)
        except Exception as e:
            raise ValueError(f"Не удалось вставить изображение: {e}")

    @staticmethod
    def _apply_font_style(font, size: Optional[int], color: Optional[List[int]]) -> None:
        if size: font.size = Pt(size)
        if color: font.color.rgb = RGBColor(*color)

    @staticmethod
    def _decode_base64_to_stream(b64_string: str) -> io.BytesIO:
        if "," in b64_string: b64_string = b64_string.split(",")[1]
        return io.BytesIO(base64.b64decode(b64_string))
